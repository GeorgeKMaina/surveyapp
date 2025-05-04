
"""
Created on Wed Mar 5, 2025
@author: Wangari Kimani
"""

import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from sklearn.feature_extraction.text import TfidfVectorizer
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from io import BytesIO
import re
import os
import joblib
from bertopic import BERTopic
from sentence_transformers import SentenceTransformer
from umap import UMAP
from hdbscan import HDBSCAN
from dotenv import load_dotenv
from openai import OpenAI
import warnings
warnings.filterwarnings("ignore", category=UserWarning, module='tqdm')

import sys
sys.path.append(r"text_utils.py")
from text_utils import clean_text

# Load environment variables
load_dotenv()

# Verify API key
api_key = os.getenv("OPENAI_API_KEY")
if not api_key:
    raise ValueError("OPENAI_API_KEY not found. Check your .env file path.")

# Initialize OpenAI client
client = OpenAI(api_key=api_key)

# Load trained question type classifier
classifier_model = joblib.load(r"best_model.pkl")

# Streamlit UI
st.title("Survey Data Analysis & PowerPoint Report Generator")

survey_type_options = ["Employee Feedback", "Customer Satisfaction", "Market Research", "Other"]
survey_type = st.selectbox("Select Survey Type", survey_type_options)
if survey_type == "Other":
    survey_type = st.text_input("Describe the survey type")

data_collection_date = st.date_input("Date of Data Collection")

uploaded_file = st.file_uploader("Upload Survey Dataset (Excel)", type=["xlsx"])

if uploaded_file:
    df = pd.ExcelFile(uploaded_file)
    sheet_name = df.sheet_names[0]
    data = df.parse(sheet_name)

    st.write("Dataset Preview:")
    st.dataframe(data.head())

    # Classify question types
    question_texts = list(data.columns)
    cleaned_questions = [clean_text(q) for q in question_texts]
    predicted_labels = classifier_model.predict(cleaned_questions)
    normalized_labels = ["open-ended" if "open" in label.lower() else "closed-ended" for label in predicted_labels]
    question_types = dict(zip(question_texts, normalized_labels))

    # Allow user to correct misclassified open-ended questions
    review_df = pd.DataFrame({
        "Question": question_texts,
        "Predicted Type": normalized_labels
    })
    st.subheader("\U0001F50D Review Question Classification")
    st.write("Below are the model's classifications. You can correct any misclassified open-ended questions:")

    corrected = st.multiselect(
        "Select questions that are actually open-ended but were misclassified:",
        options=review_df[review_df["Predicted Type"] == "closed-ended"]["Question"].tolist()
    )

    for q in corrected:
        question_types[q] = "open-ended"

    open_ended = [q for q, label in question_types.items() if label == 'open-ended']
    closed_ended = [q for q, label in question_types.items() if label == 'closed-ended']

    # Show final classification table
    final_df = pd.DataFrame({
        "Question": question_texts,
        "Final Type": [question_types[q] for q in question_texts]
    })
    st.subheader("\U0001F4BE Final Classification Results")
    st.dataframe(final_df)

    def preprocess_text(text):
        if isinstance(text, str):
            text = text.lower()
            text = re.sub(r'[^a-zA-Z\s]', '', text)
            return text
        return ""

    data_cleaned = data[open_ended].astype(str).applymap(preprocess_text)
    topic_freqs = {}
    closed_insights = {}

    for question in open_ended:
        responses = data_cleaned[question].dropna().tolist()
        responses = [r for r in responses if isinstance(r, str) and len(r.strip()) > 3]

        if len(responses) > 5:
            try:
                embedding_model = SentenceTransformer("all-MiniLM-L6-v2")
                umap_model = UMAP(n_neighbors=15, n_components=5, min_dist=0.1, random_state=42)
                hdbscan_model = HDBSCAN(min_cluster_size=5, prediction_data=True)
                vectorizer_model = TfidfVectorizer(max_features=1000, min_df=1, max_df=1.0, stop_words="english")

                dynamic_topic_model = BERTopic(
                    embedding_model=embedding_model,
                    umap_model=umap_model,
                    hdbscan_model=hdbscan_model,
                    vectorizer_model=vectorizer_model,
                    nr_topics=8,
                    calculate_probabilities=False,
                    verbose=False
                )

                topics, _ = dynamic_topic_model.fit_transform(responses)

                topic_counts = pd.Series(topics).value_counts().sort_values(ascending=False)
                filtered = {
                    " / ".join([word for word, _ in dynamic_topic_model.get_topic(topic)[:3]]): count
                    for topic, count in topic_counts.items()
                    if topic != -1 and count >= 3
                }
                topic_freqs[question] = filtered if filtered else {"No dominant topics": 1}
            except Exception as e:
                topic_freqs[question] = {f"BERTopic Error: {str(e)}": 1}
        else:
            topic_freqs[question] = {"Not enough responses": 1}

    for question in closed_ended:
        responses = data[question].value_counts(normalize=True) * 100
        closed_insights[question] = responses.to_dict()

    # PowerPoint generation
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Executive Summary"

    executive_prompt = f"""
    You are a senior data analyst summarizing survey findings.
    Survey type: {survey_type}, conducted on {data_collection_date.strftime('%Y-%m-%d')}.

    Closed-ended question summaries:
    {closed_insights}

    Open-ended topic highlights:
    {topic_freqs}

    Write a concise executive summary (5–7 bullet points) highlighting:
    - Key patterns in closed responses
    - Repeated themes in open-ended answers
    - Strategic takeaways or actions
    """

    try:
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": executive_prompt}]
        )
        summary_text = response.choices[0].message.content
    except Exception as e:
        summary_text = f"Executive summary could not be generated. Error: {e}"

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5.5))
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.text = summary_text
    text_frame.word_wrap = True
    text_frame.auto_size = True

    for question in closed_ended:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = question
        responses = data[question].value_counts(normalize=True) * 100
        if responses.empty:
            responses = pd.Series([0], index=['No Data'])

        chart_data = CategoryChartData()
        chart_data.categories = list(responses.index)
        chart_data.add_series("Responses", tuple(responses.values))

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(10.5), Inches(3), chart_data
        )
        chart.chart.has_title = False

        prompt = f"""
        You are a data analyst preparing a summary for this survey question:
        "{question}"

        Response distribution: {responses.to_dict()}

        Write 3–5 bullet points interpreting these results and suggesting next steps.
        """

        try:
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}]
            )
            narrative = response.choices[0].message.content
        except Exception as e:
            narrative = f"Error generating summary: {e}"

        textbox = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(10.5), Inches(2.5))
        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.text = narrative
        text_frame.word_wrap = True
        text_frame.auto_size = True

    for question, topics in topic_freqs.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = question

        if topics and isinstance(topics, dict):
            categories = list(topics.keys())
            values = list(topics.values())

            if categories:
                chart_data = CategoryChartData()
                chart_data.categories = categories
                chart_data.add_series("Mentions", values)

                chart = slide.shapes.add_chart(
                    XL_CHART_TYPE.BAR_CLUSTERED, Inches(1), Inches(1), Inches(10.5), Inches(3), chart_data
                )
                chart.chart.has_title = False

                prompt = f"""
                You are a data analyst summarizing open-ended survey feedback.

                Question: "{question}"

                Top discussed topics and their counts:
                {dict(zip(categories, values))}

                Write 3–5 concise bullet points summarizing:
                - The main themes
                - What participants are emphasizing
                - Any suggested actions based on feedback
                """

                try:
                    response = client.chat.completions.create(
                        model="gpt-3.5-turbo",
                        messages=[{"role": "user", "content": prompt}]
                    )
                    narrative = response.choices[0].message.content
                except Exception as e:
                    narrative = f"Error generating summary: {e}"

                textbox = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(10.5), Inches(2.5))
                text_frame = textbox.text_frame
                text_frame.clear()
                text_frame.text = narrative
                text_frame.word_wrap = True
                text_frame.auto_size = True
                continue

        tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10.5), Inches(4))
        text_frame = tb.text_frame
        text_frame.clear()
        text_frame.text = "No topics generated or too few responses."
        text_frame.word_wrap = True
        text_frame.auto_size = True

    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)

    st.success("\u2705 Analysis complete! Download your PowerPoint report below.")
    st.download_button(
        "\U0001F4E5 Download PowerPoint Report",
        ppt_stream,
        "survey_report.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
